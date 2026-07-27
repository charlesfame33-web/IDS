# -*- coding: utf-8 -*-
"""Generate 12-slide BSc project defence slide decks for both projects,
using the standard CSC / OAUSTECH format.  Black & white throughout."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


def build_deck(cfg, out_path):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    DARK = RGBColor(*cfg['dark'])
    PANEL = RGBColor(*cfg['panel'])
    ACCENT = RGBColor(*cfg['accent'])
    ACCENT2 = RGBColor(*cfg['accent2'])
    WHITE = RGBColor(*cfg.get('text', (0xF5, 0xF7, 0xFA)))
    MUTED = RGBColor(*cfg.get('muted', (0x9A, 0xA5, 0xB1)))
    BORDER = RGBColor(*cfg.get('border', (0x30, 0x36, 0x3D)))

    def bg(slide, color=DARK):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def box(slide, l, t, w, h, color, border=None):
        sp = slide.shapes.add_shape(1, l, t, w, h)
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        if border is not None:
            sp.line.color.rgb = border; sp.line.width = Pt(1)
        else:
            sp.line.fill.background()
        sp.shadow.inherit = False
        return sp

    def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             space=6):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = anchor
        for i, (txt, size, color, bold) in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align; p.space_after = Pt(space)
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = 'Segoe UI'
        return tb

    def bullets(slide, l, t, w, h, items, size=18, color=WHITE, gap=10):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap)
            r = p.add_run(); r.text = '\u25B8  ' + it
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.name = 'Segoe UI'
        return tb

    def accent_bar(slide, t=Inches(1.15)):
        box(slide, Inches(0.7), t, Inches(1.6), Pt(5), ACCENT)

    def header(slide, kicker, title):
        text(slide, Inches(0.7), Inches(0.45), Inches(12), Inches(0.4),
             [(kicker.upper(), 14, ACCENT, True)])
        text(slide, Inches(0.7), Inches(0.75), Inches(12), Inches(0.9),
             [(title, 30, WHITE, True)])
        accent_bar(slide, Inches(1.5))

    # ---------- Title slide ----------
    s = prs.slides.add_slide(blank); bg(s)
    box(s, 0, 0, W, Inches(0.18), ACCENT)
    box(s, Inches(0.7), Inches(1.6), Inches(0.14), Inches(3.4), ACCENT)
    text(s, Inches(1.05), Inches(1.5), Inches(11.5), Inches(0.5),
         [(cfg['badge'].upper(), 15, ACCENT2, True)])
    text(s, Inches(1.05), Inches(2.0), Inches(11.6), Inches(2.4),
         [(cfg['title'], 40, WHITE, True)])
    text(s, Inches(1.05), Inches(4.7), Inches(11.5), Inches(1.6),
         [(f'Presented by {cfg["student"]}', 18, WHITE, False),
          (f'Matric No: {cfg["matric"]}', 16, MUTED, False),
          (f'Supervisor: {cfg["supervisor"]}', 16, MUTED, False),
          ('Dept. of Computer Science \u2014 OAUSTECH, Okitipupa', 16, MUTED, False)])
    box(s, 0, H - Inches(0.18), W, Inches(0.18), ACCENT2)

    # ---------- Content slides ----------
    for sl in cfg['slides']:
        s = prs.slides.add_slide(blank); bg(s)
        header(s, sl['kicker'], sl['title'])
        kind = sl.get('kind', 'bullets')
        if kind == 'bullets':
            bullets(s, Inches(0.75), Inches(1.9), Inches(11.9), Inches(5.2),
                    sl['items'], size=sl.get('size', 18))
        elif kind == 'two':
            box(s, Inches(0.75), Inches(1.9), Inches(5.75), Inches(4.9), PANEL, border=BORDER)
            box(s, Inches(6.8), Inches(1.9), Inches(5.75), Inches(4.9), PANEL, border=BORDER)
            text(s, Inches(1.0), Inches(2.1), Inches(5.3), Inches(0.5),
                 [(sl['left_title'], 19, ACCENT, True)])
            bullets(s, Inches(1.0), Inches(2.75), Inches(5.3), Inches(3.9),
                    sl['left'], size=15, gap=7)
            text(s, Inches(7.05), Inches(2.1), Inches(5.3), Inches(0.5),
                 [(sl['right_title'], 19, ACCENT2, True)])
            bullets(s, Inches(7.05), Inches(2.75), Inches(5.3), Inches(3.9),
                    sl['right'], size=15, gap=7)
        elif kind == 'cards':
            cards = sl['cards']; n = len(cards)
            cols = min(3, max(2, (n + 1) // 2))
            import math
            rows = math.ceil(n / cols)
            gap = Inches(0.3)
            cw = (Inches(11.9) - gap * (cols - 1)) / cols
            ch = (Inches(4.9) - gap * (rows - 1)) / rows
            for idx, (ct, cb) in enumerate(cards):
                r, c = idx // cols, idx % cols
                l = Inches(0.75) + c * (cw + gap)
                t = Inches(1.95) + r * (ch + gap)
                box(s, l, t, cw, ch, PANEL, border=BORDER)
                box(s, l, t, Pt(5), ch, ACCENT if idx % 2 == 0 else ACCENT2)
                text(s, l + Inches(0.25), t + Inches(0.18), cw - Inches(0.5),
                     Inches(0.5), [(ct, 17, WHITE, True)])
                text(s, l + Inches(0.25), t + Inches(0.75), cw - Inches(0.5),
                     ch - Inches(0.9), [(cb, 13, MUTED, False)])
        elif kind == 'implementation':
            items = sl['items']
            ft = Inches(1.95)
            fw = Inches(11.9)
            fh = Inches(3.2)
            box(s, Inches(0.75), ft, fw, fh, PANEL, border=BORDER)
            text(s, Inches(1.05), ft + Inches(0.15), fw - Inches(0.6),
                 Inches(0.35), [(items[0][0], 15, ACCENT, True)])
            if items[0][2] and os.path.exists(items[0][2]):
                try:
                    s.shapes.add_picture(items[0][2], Inches(1.1), ft + Inches(0.55),
                                         Inches(11.2))
                except Exception:
                    pass
            st = ft + fh + Inches(0.25)
            for idx in range(1, len(items)):
                fig_label, fig_desc, img_path = items[idx]
                l = Inches(0.75) + (idx - 1) * Inches(6.05)
                cw = Inches(5.75)
                ch = Inches(1.8)
                box(s, l, st, cw, ch, PANEL, border=BORDER)
                text(s, l + Inches(0.25), st + Inches(0.12), cw - Inches(0.5),
                     Inches(0.3), [(fig_label, 13, ACCENT, True)])
                if img_path and os.path.exists(img_path):
                    try:
                        s.shapes.add_picture(img_path, l + Inches(0.3), st + Inches(0.5),
                                             cw - Inches(0.6))
                    except Exception:
                        pass
                text(s, l + Inches(0.25), st + ch - Inches(0.35),
                     cw - Inches(0.5), Inches(0.3),
                     [(fig_desc, 11, MUTED, False)], align=PP_ALIGN.CENTER)
        elif kind == 'screenshot':
            items = sl['items']
            for idx, item in enumerate(items):
                fig_label, fig_desc, img_path = item if len(item) == 3 else (item[0], item[1], None)
                l = Inches(0.75) + idx * (Inches(6.05))
                t = Inches(1.95)
                cw = Inches(5.75)
                ch = Inches(4.85)
                box(s, l, t, cw, ch, PANEL, border=BORDER)
                text(s, l + Inches(0.3), t + Inches(0.2), cw - Inches(0.6),
                     Inches(0.4), [(fig_label, 15, ACCENT, True)])
                if img_path and os.path.exists(img_path):
                    try:
                        s.shapes.add_picture(img_path, l + Inches(0.35),
                                             t + Inches(0.72), cw - Inches(0.7))
                    except Exception:
                        pass
                else:
                    inner = box(s, l + Inches(0.3), t + Inches(0.7),
                                cw - Inches(0.6), ch - Inches(1.2),
                                RGBColor(0xE8, 0xE8, 0xE8), border=BORDER)
                    text(s, l + Inches(0.3), t + ch / 2 - Inches(0.3),
                         cw - Inches(0.6), Inches(0.5),
                         [('Screenshot' if not img_path else f'Missing: {os.path.basename(img_path)}',
                           13, RGBColor(0x99, 0x99, 0x99), False)],
                         align=PP_ALIGN.CENTER)
                text(s, l + Inches(0.3), t + ch - Inches(0.4),
                     cw - Inches(0.6), Inches(0.35),
                     [(fig_desc, 11.5, MUTED, False)],
                     align=PP_ALIGN.CENTER)

    # ---------- Thank-you ----------
    s = prs.slides.add_slide(blank); bg(s)
    box(s, 0, 0, W, Inches(0.18), ACCENT)
    text(s, Inches(0.7), Inches(2.7), Inches(12), Inches(1.2),
         [('Thank You', 48, WHITE, True)], align=PP_ALIGN.CENTER)
    text(s, Inches(0.7), Inches(3.9), Inches(12), Inches(0.8),
         [('Questions & Demonstration', 22, ACCENT, True)], align=PP_ALIGN.CENTER)
    box(s, 0, H - Inches(0.18), W, Inches(0.18), ACCENT2)

    prs.save(out_path)
    print('saved', out_path, '-', len(prs.slides.__iter__.__self__._sldIdLst), 'slides')


# ==================== BLACK & WHITE THEME ====================
BW = {
    'dark': (0xFF, 0xFF, 0xFF),
    'panel': (0xF4, 0xF4, 0xF4),
    'accent': (0x00, 0x00, 0x00),
    'accent2': (0x44, 0x44, 0x44),
    'text': (0x00, 0x00, 0x00),
    'muted': (0x33, 0x33, 0x33),
    'border': (0x00, 0x00, 0x00),
}

# ==================== MAYOR4CODE DECK ====================
mayor = {
    **BW,
    'badge': 'BSc Project Defence',
    'title': 'Design and Implementation of a Web-Based Interactive '
             'E-Learning Platform for Structured Python Programming',
    'student': 'Giwa Mayowa Bopoola',
    'matric': 'CSC/22/204',
    'supervisor': "'Dr. Adeolu Obamehinti'",
    'slides': [
        # --- Slide 2: Introduction ---
        {'kicker': 'Slide 2', 'title': 'Introduction',
         'items': [
             'Self-paced, learner-centred education (Garrison, 2017).',
             'Interactive platforms improve engagement through instant feedback (Pressman & Maxim, 2020).',
             'Python is widely used, but beginners struggle with traditional methods.',
             'mayor4code: integrated platform for structured Python learning.',
         ]},
        # --- Slide 3: Motivation / Problem Statement ---
        {'kicker': 'Slide 3', 'title': 'Motivation & Problem Statement',
         'items': [
             'Unstructured content confuses learners on what to study next.',
             'Theory and practice separated; local setup is discouraging.',
             'Few resources provide instant feedback or enforce mastery.',
             'No progress tracking, certification, or motivation features.',
         ]},
        # --- Slide 4: Aim & Objectives ---
        {'kicker': 'Slide 4', 'title': 'Aim and Objectives',
         'kind': 'two',
         'left_title': 'Aim',
         'left': ['Design & implement a web-based e-learning platform integrating '
                  'lessons, coding, and assessment in one system.'],
         'right_title': 'Objectives',
         'right': ['Structured lessons with locked progression',
                   'Secure user authentication',
                   'In-browser Python playground',
                   'Automated quiz scoring',
                   'Progress tracking, certification, leaderboard',
                   'Evaluate through testing']},
        # --- Slide 5: Significance ---
        {'kicker': 'Slide 5', 'title': 'Significance of the Study',
         'items': [
             'Educational: structured, self-paced learning with instant feedback.',
             'Technological: modern web stack + secure code execution.',
             'Economic: free platform, no per-learner cost.',
             'Social: accessible to anyone with internet.',
         ]},
        # --- Slide 6: Literature Review ---
        {'kicker': 'Slide 6', 'title': 'Literature Review',
         'items': [
             'E-learning as effective as traditional instruction (Anderson, 2008).',
             'Constructivism supports interactive, scaffolded learning (Piaget, 1952; Vygotsky, 1978).',
             'Mastery learning validates locked progression (Bloom, 1968; Guskey, 2007).',
             'Gamification boosts engagement (Deterding et al., 2011; Hamari et al., 2014).',
             'Gap: no free platform integrates all features in one system.',
         ]},
        # --- Slide 7: Methodology ---
        {'kicker': 'Slide 7', 'title': 'Methodology',
         'kind': 'cards',
         'cards': [
             ('Approach', 'Agile SDLC with iterative testing.'),
             ('Architecture', 'Three-tier: Presentation, Django, Database.'),
             ('Stack', 'Django, Python, HTML/CSS/JS; SQLite/PostgreSQL.'),
             ('Progression', 'Lessons unlock at \u226560% quiz score.'),
             ('Playground', 'Isolated subprocess with 5s timeout.'),
             ('Testing', 'Unit, functional, and usability testing.')]},
        # --- Slide 8: Implementation ---
        {'kicker': 'Slide 8', 'title': 'Implementation',
         'kind': 'implementation',
         'items': [
             ('System Flowchart',
              'User journey: register \u2192 learn \u2192 quiz \u2192 unlock \u2192 certificate.',
              r'c:\Users\ALEXIS\Desktop\SENPAI\shots\mayor4code\flowchart.png'),
             ('Lesson Interface',
              'Structured Python lessons with navigation.',
              r'c:\Users\ALEXIS\Desktop\SENPAI\shots\mayor4code\06-lessons.png'),
             ('Code Playground',
              'In-browser editor with safe execution.',
              r'c:\Users\ALEXIS\Desktop\SENPAI\shots\mayor4code\07-playground.png'),
         ]},
        # --- Slide 9: Results & Discussion ---
        {'kicker': 'Slide 9', 'title': 'Results & Discussion',
         'kind': 'screenshot',
         'items': [
             ('Quiz Interface',
              'MCQ with instant scoring and pass-mark feedback.',
              r'c:\Users\ALEXIS\Desktop\SENPAI\shots\mayor4code\08-quizzes.png'),
             ('Learner Dashboard',
              'Progress bar and lesson unlock status.',
              r'c:\Users\ALEXIS\Desktop\SENPAI\shots\mayor4code\05-dashboard.png'),
         ]},
        # --- Slide 10: Contributions to Knowledge ---
        {'kicker': 'Slide 10', 'title': 'Contributions to Knowledge',
         'items': [
             'Integrated platform: lessons, playground, assessment, gamification.',
             'Mastery-based locked progression model.',
             'Secure in-browser code execution via isolated subprocess.',
             'Automated certification with unique verification codes.',
         ]},
        # --- Slide 11: Conclusion & References ---
        {'kicker': 'Slide 11', 'title': 'Conclusion & References',
         'kind': 'two',
         'left_title': 'Conclusion',
         'left': ['All 6 objectives achieved: structured lessons, secure auth, '
                  'playground, auto-scoring, certification, leaderboard.',
                  'Testing confirmed all features work correctly.',
                  'Integrated platform effective for Python beginners.'],
         'right_title': 'References',
         'right': ['Anderson, T. (2008). The theory and practice of online learning.',
                   'Bloom, B. S. (1968). Mastery learning. Evaluation Comment, 1(2).',
                   'Garrison, D. R. (2017). E-learning in the 21st century (3rd ed.). Routledge.',
                   'Pressman, R. S., & Maxim, B. R. (2020). Software engineering (9th ed.). McGraw-Hill.',
                   'Vygotsky, L. S. (1978). Mind in society. Harvard University Press.']},
    ],
}

# ==================== AGRIFLOW DECK (cream + green + black) ====================
AGRI_WARM = {
    'dark': (0xF8, 0xF3, 0xEE),
    'panel': (0xF0, 0xEB, 0xE6),
    'accent': (0x05, 0x96, 0x69),
    'accent2': (0x04, 0x7A, 0x55),
    'text': (0x00, 0x00, 0x00),
    'muted': (0x33, 0x33, 0x33),
    'border': (0x05, 0x96, 0x69),
}
agri = {
    **AGRI_WARM,
    'badge': 'BSc Project Defence',
    'title': 'AI-Driven Agricultural Supply Chain & Produce '
             'Scheduling System',
    'student': 'Obayomi Samuel Oluwagbotemi',
    'matric': 'CSC/22/124',
    'supervisor': "'Dr. Adeolu Obamehinti'",
    'slides': [
        # --- Slide 2: Introduction ---
        {'kicker': 'Slide 2', 'title': 'Introduction',
         'items': [
             'Agricultural supply chains for fresh produce suffer significant post-harvest '
             'losses due to poor coordination, inefficient logistics, and lack of visibility.',
             'Farmers, buyers, transporters, and warehouse managers often operate in '
             'isolation, relying on manual and informal processes.',
             'Web-based platforms and data-driven techniques can improve coordination '
             'and reduce waste (Pressman & Maxim, 2020).',
             'AgriFlow AI is an intelligent web platform that unifies supply chain '
             'participants and applies deterministic decision rules to improve efficiency.',
             'The system integrates a marketplace, perishable-first scheduling, demand '
             'forecasting, spoilage tracking, and an AI assistant grounded in live data.',
         ]},
        # --- Slide 3: Motivation / Problem Statement ---
        {'kicker': 'Slide 3', 'title': 'Motivation / Statement of the Problem',
         'items': [
             'A substantial proportion of fresh produce is lost between harvest and '
             'consumption due to poor logistics and weak coordination.',
             'Deliveries are frequently scheduled without regard to perishability; '
             'the most time-critical produce is not prioritised.',
             'Demand is rarely anticipated systematically, causing both shortages '
             'and surpluses that increase waste.',
             'Operational data is often complex and difficult for non-experts to interpret, '
             'limiting its usefulness for decision-making.',
             'These challenges motivated the development of AgriFlow AI as an integrated, '
             'intelligent, and practical platform for produce supply chain management.',
         ]},
        # --- Slide 4: Aim & Objectives ---
        {'kicker': 'Slide 4', 'title': 'Aim and Objectives of the Study',
         'kind': 'two',
         'left_title': 'Aim',
         'left': ['To design and implement an intelligent web-based agricultural supply '
                  'chain platform that unifies participants and applies data-driven decision '
                  'rules to improve efficiency and reduce post-harvest losses.'],
         'right_title': 'Objectives',
         'right': ['Design a role-based platform for 5 user types',
                   'Implement a produce marketplace with order management',
                   'Develop a deterministic perishable-first scheduling engine',
                   'Implement demand forecasting and shelf-life spoilage tracking',
                   'Integrate an AI assistant grounded in live operational data',
                   'Evaluate functionality and reliability through testing']},
        # --- Slide 5: Significance ---
        {'kicker': 'Slide 5', 'title': 'Significance of the Study',
         'items': [
             'Economic Impact: Reduces post-harvest losses through perishable-first '
             'scheduling and spoilage tracking, improving incomes and lowering costs.',
             'Social Impact: Strengthens coordination among supply chain participants '
             'and contributes to food security by reducing food waste.',
             'Technological Impact: Demonstrates integration of web technologies, '
             'deterministic optimisation, and grounded AI for agricultural challenges.',
             'Practical Impact: Provides a working model for combining reliable rule-based '
             'logic with accessible natural-language explanations for non-expert users.',
         ]},
        # --- Slide 6: Literature Review ---
        {'kicker': 'Slide 6', 'title': 'Literature Review',
         'items': [
             'Post-harvest loss is a central concern in agricultural supply chains; effective '
             'management minimises losses through improved coordination (Christopher, 2016).',
             'Electronic marketplaces improve access but do not solve logistical challenges; '
             'integrated platforms deliver greater value (Sommerville, 2016).',
             'Deterministic rules such as perishable-first scheduling, moving averages, '
             'and haversine distance give reliable, explainable decisions.',
             'AI must be grounded in verified data to avoid fabrication; it should '
             'explain rather than decide (Russell & Norvig, 2021).',
             'Gap: Few systems integrate marketplace, perishability-aware logistics, '
             'forecasting, and grounded AI assistance in one platform.',
         ]},
        # --- Slide 7: Methodology ---
        {'kicker': 'Slide 7', 'title': 'Methodology',
         'kind': 'cards',
         'cards': [
             ('Approach', 'Agile SDLC with automated end-to-end testing across all roles.'),
             ('Architecture', 'Serverless: Next.js server components + cloud database.'),
             ('Stack', 'Next.js 14, TypeScript, Tailwind; Supabase PostgreSQL + RLS.'),
             ('Scheduling', 'Perishable-first; nearest warehouse; least-busy transporter.'),
             ('Forecasting', 'Simple & weighted moving averages over 8 weeks.'),
             ('AI Assistant', 'Google Gemini \u2014 explains live data only; all '
                              'decisions are deterministic.')]},
        # --- Slide 8: Implementation (screenshots) ---
        {'kicker': 'Slide 8', 'title': 'Implementation',
         'kind': 'screenshot',
         'items': [
              ('Figure 1: Farmer & Buyer Dashboards',
               'Role-based dashboards for crop management, marketplace browsing, and orders.',
               r'c:\Users\ALEXIS\Desktop\SENPAI\shots\agriflow\dash-farmer.png'),
              ('Figure 2: Scheduling & AI Assistant',
               'Perishable-first scheduling engine and AI assistant grounded in live data.',
               r'c:\Users\ALEXIS\Desktop\SENPAI\shots\agriflow\dash-admin.png'),
         ]},
        # --- Slide 9: Results & Discussion (screenshots) ---
        {'kicker': 'Slide 9', 'title': 'Results & Discussion',
         'kind': 'screenshot',
         'items': [
              ('Figure 3: Warehouse Inventory',
               'Inventory management with live shelf-life alerts and spoilage tracking.',
               r'c:\Users\ALEXIS\Desktop\SENPAI\shots\agriflow\warehouse-inventory.png'),
              ('Figure 4: Transport & Admin Views',
               'Delivery workflow, route map, and admin scheduling dashboard.',
               r'c:\Users\ALEXIS\Desktop\SENPAI\shots\agriflow\dash-transporter.png'),
         ]},
        # --- Slide 10: Contributions ---
        {'kicker': 'Slide 10', 'title': 'Contributions to Knowledge',
         'items': [
             'A unified, role-based platform integrating marketplace, logistics, inventory, '
             'and analytics for agricultural supply chain management.',
             'A deterministic perishable-first scheduling engine that prioritises deliveries '
             'by remaining shelf life and optimises warehouse and transporter assignments.',
             'Demand forecasting using simple and weighted moving averages, with trend '
             'indicators for accessible planning support.',
             'Continuous shelf-life spoilage tracking that classifies stored items as '
             'fresh, expiring, or spoiled for proactive decision-making.',
             'An AI assistant grounded in live operational data that provides '
             'natural-language explanations, making complex data accessible to all users.',
         ]},
        # --- Slide 11: Conclusion + References ---
        {'kicker': 'Slide 11', 'title': 'Conclusion / References',
         'kind': 'two',
         'left_title': 'Conclusion',
         'left': ['All objectives were achieved: a role-based platform for five user types, '
                  'a marketplace with order management, a perishable-first scheduling '
                  'engine, demand forecasting, spoilage tracking, and an explanatory '
                  'AI assistant were implemented and tested.',
                  'The platform demonstrates that combining a unified digital system with '
                  'deterministic optimisation and grounded AI improves supply chain efficiency.'],
         'right_title': 'References',
         'right': ['Christopher, M. (2016). Logistics and supply chain management (5th ed.). '
                   'Pearson.',
                   'Elmasri, R. & Navathe, S. B. (2017). Fundamentals of database systems '
                   '(7th ed.). Pearson.',
                   'Google. (2024). Gemini API documentation. ai.google.dev',
                   'Pressman, R. S. & Maxim, B. R. (2020). Software engineering (9th ed.). '
                   'McGraw-Hill.',
                   'Russell, S. J. & Norvig, P. (2021). Artificial intelligence: A modern '
                   'approach (4th ed.). Pearson.',
                   'Sommerville, I. (2016). Software engineering (10th ed.). Pearson.',
                   'Supabase. (2024). Supabase documentation. supabase.com/docs']},
    ],
}

# ==================== AGRIFLOW DECK (plain — no screenshots) ====================
agri_plain = {
    **AGRI_WARM,
    'badge': 'BSc Project Defence',
    'title': 'AI-Driven Agricultural Supply Chain & Produce '
             'Scheduling System',
    'student': 'Obayomi Samuel Oluwagbotemi',
    'matric': 'CSC/22/124',
    'supervisor': "'Dr. Adeolu Obamehinti'",
    'slides': [
        {'kicker': 'Slide 2', 'title': 'Introduction',
         'items': [
             'Agricultural supply chains for fresh produce suffer significant post-harvest '
             'losses due to poor coordination, inefficient logistics, and lack of visibility.',
             'Farmers, buyers, transporters, and warehouse managers often operate in '
             'isolation, relying on manual and informal processes.',
             'Web-based platforms and data-driven techniques can improve coordination '
             'and reduce waste (Pressman & Maxim, 2020).',
             'AgriFlow AI is an intelligent web platform that unifies supply chain '
             'participants and applies deterministic decision rules to improve efficiency.',
             'The system integrates a marketplace, perishable-first scheduling, demand '
             'forecasting, spoilage tracking, and an AI assistant grounded in live data.',
         ]},
        {'kicker': 'Slide 3', 'title': 'Motivation / Statement of the Problem',
         'items': [
             'A substantial proportion of fresh produce is lost between harvest and '
             'consumption due to poor logistics and weak coordination.',
             'Deliveries are frequently scheduled without regard to perishability; '
             'the most time-critical produce is not prioritised.',
             'Demand is rarely anticipated systematically, causing both shortages '
             'and surpluses that increase waste.',
             'Operational data is often complex and difficult for non-experts to interpret, '
             'limiting its usefulness for decision-making.',
             'These challenges motivated the development of AgriFlow AI as an integrated, '
             'intelligent, and practical platform for produce supply chain management.',
         ]},
        {'kicker': 'Slide 4', 'title': 'Aim and Objectives of the Study',
         'kind': 'two',
         'left_title': 'Aim',
         'left': ['To design and implement an intelligent web-based agricultural supply '
                  'chain platform that unifies participants and applies data-driven decision '
                  'rules to improve efficiency and reduce post-harvest losses.'],
         'right_title': 'Objectives',
         'right': ['Design a role-based platform for 5 user types',
                   'Implement a produce marketplace with order management',
                   'Develop a deterministic perishable-first scheduling engine',
                   'Implement demand forecasting and shelf-life spoilage tracking',
                   'Integrate an AI assistant grounded in live operational data',
                   'Evaluate functionality and reliability through testing']},
        {'kicker': 'Slide 5', 'title': 'Significance of the Study',
         'items': [
             'Economic Impact: Reduces post-harvest losses through perishable-first '
             'scheduling and spoilage tracking, improving incomes and lowering costs.',
             'Social Impact: Strengthens coordination among supply chain participants '
             'and contributes to food security by reducing food waste.',
             'Technological Impact: Demonstrates integration of web technologies, '
             'deterministic optimisation, and grounded AI for agricultural challenges.',
             'Practical Impact: Provides a working model for combining reliable rule-based '
             'logic with accessible natural-language explanations for non-expert users.',
         ]},
        {'kicker': 'Slide 6', 'title': 'Literature Review',
         'items': [
             'Post-harvest loss is a central concern in agricultural supply chains; effective '
             'management minimises losses through improved coordination (Christopher, 2016).',
             'Electronic marketplaces improve access but do not solve logistical challenges; '
             'integrated platforms deliver greater value (Sommerville, 2016).',
             'Deterministic rules such as perishable-first scheduling, moving averages, '
             'and haversine distance give reliable, explainable decisions.',
             'AI must be grounded in verified data to avoid fabrication; it should '
             'explain rather than decide (Russell & Norvig, 2021).',
             'Gap: Few systems integrate marketplace, perishability-aware logistics, '
             'forecasting, and grounded AI assistance in one platform.',
         ]},
        {'kicker': 'Slide 7', 'title': 'Methodology',
         'kind': 'cards',
         'cards': [
             ('Approach', 'Agile SDLC with automated end-to-end testing across all roles.'),
             ('Architecture', 'Serverless: Next.js server components + cloud database.'),
             ('Stack', 'Next.js 14, TypeScript, Tailwind; Supabase PostgreSQL + RLS.'),
             ('Scheduling', 'Perishable-first; nearest warehouse; least-busy transporter.'),
             ('Forecasting', 'Simple & weighted moving averages over 8 weeks.'),
             ('AI Assistant', 'Google Gemini \u2014 explains live data only; all '
                              'decisions are deterministic.')]},
        {'kicker': 'Slide 8', 'title': 'Implementation',
         'kind': 'bullets',
         'items': [
             'Role-based platform built with Next.js 14, TypeScript, and Tailwind CSS.',
             'Supabase PostgreSQL with row-level security for data access control.',
             'Deterministic algorithms: perishable-first scheduling, moving-average forecasting.',
             'Google Gemini AI assistant grounded in live operational data.',
             'Automated end-to-end testing across all five roles with Playwright.',
             'Deployed as a progressive web application on Vercel.']},
        {'kicker': 'Slide 9', 'title': 'Results & Discussion',
         'kind': 'bullets',
         'items': [
             'All functional and non-functional requirements were met.',
             '21 test cases across unit, functional, and E2E levels all passed.',
             'Perishable-first scheduling correctly prioritises shortest shelf life first.',
             'Demand forecasting computes accurate moving averages with trend indicators.',
             'Shelf-life spoilage tracking correctly classifies items as fresh/expiring/spoiled.',
             'AI assistant provides accurate plain-language explanations of operational data.',
         ]},
        {'kicker': 'Slide 10', 'title': 'Contributions to Knowledge',
         'items': [
             'A unified, role-based platform integrating marketplace, logistics, inventory, '
             'and analytics for agricultural supply chain management.',
             'A deterministic perishable-first scheduling engine that prioritises deliveries '
             'by remaining shelf life and optimises warehouse and transporter assignments.',
             'Demand forecasting using simple and weighted moving averages, with trend '
             'indicators for accessible planning support.',
             'Continuous shelf-life spoilage tracking that classifies stored items as '
             'fresh, expiring, or spoiled for proactive decision-making.',
             'An AI assistant grounded in live operational data that provides '
             'natural-language explanations, making complex data accessible to all users.',
         ]},
        {'kicker': 'Slide 11', 'title': 'Conclusion / References',
         'kind': 'two',
         'left_title': 'Conclusion',
         'left': ['All objectives were achieved: a role-based platform for five user types, '
                  'a marketplace with order management, a perishable-first scheduling '
                  'engine, demand forecasting, spoilage tracking, and an explanatory '
                  'AI assistant were implemented and tested.',
                  'The platform demonstrates that combining a unified digital system with '
                  'deterministic optimisation and grounded AI improves supply chain efficiency.'],
         'right_title': 'References',
         'right': ['Christopher, M. (2016). Logistics and supply chain management (5th ed.). '
                   'Pearson.',
                   'Elmasri, R. & Navathe, S. B. (2017). Fundamentals of database systems '
                   '(7th ed.). Pearson.',
                   'Google. (2024). Gemini API documentation. ai.google.dev',
                   'Pressman, R. S. & Maxim, B. R. (2020). Software engineering (9th ed.). '
                   'McGraw-Hill.',
                   'Russell, S. J. & Norvig, P. (2021). Artificial intelligence: A modern '
                   'approach (4th ed.). Pearson.',
                   'Sommerville, I. (2016). Software engineering (10th ed.). Pearson.',
                   'Supabase. (2024). Supabase documentation. supabase.com/docs']},
    ],
}

# ==================== IDS DECK ====================
ids = {
    **BW,
    'badge': 'BSc Project Defence',
    'title': 'Design and Implementation of an AI-Based Intrusion Detection '
             'System for Encrypted Traffic',
    'student': 'Oyeduntan Segun Elijah',
    'matric': 'CSC/22/174',
    'supervisor': "'Dr. (Engr.) Modupe Agagu'",
    'slides': [
        {'kicker': 'Slide 2', 'title': 'Introduction',
         'items': [
             'Network traffic is increasingly encrypted, creating a gap where '
             'traditional signature-based intrusion detection systems cannot inspect payloads.',
             'Machine learning models can detect malicious patterns in flow-based features '
             'without decrypting traffic, preserving privacy while maintaining security.',
             'XGBoost, a gradient-boosting ensemble, achieves state-of-the-art results on '
             'tabular network data with high accuracy and low inference latency.',
             'This project presents an AI-based IDS that combines an XGBoost model trained '
             'on the CICIDS2017 dataset with an interactive Streamlit dashboard.',
             'The system supports CSV upload, PCAP file analysis, and live network capture, '
             'with an optional Gemini AI assistant for plain-English explanations.',
         ]},
        {'kicker': 'Slide 3', 'title': 'Motivation / Statement of the Problem',
         'items': [
             'Encrypted traffic now accounts for over 90% of internet traffic, making '
             'deep-packet-inspection-based IDS largely ineffective.',
             'Signature-based systems cannot detect novel or zero-day attacks, which '
             'increasingly exploit encrypted channels.',
             'Existing ML-based IDS solutions often lack user-friendly interfaces, making '
             'them inaccessible to non-experts.',
             'There is a need for a system that combines accurate ML detection, support '
             'for multiple input modes, and an intuitive interface.',
             'These challenges motivated the development of an accessible, AI-powered '
             'IDS that works with encrypted traffic.',
         ]},
        {'kicker': 'Slide 4', 'title': 'Aim and Objectives of the Study',
         'kind': 'two',
         'left_title': 'Aim',
         'left': ['To design and implement an AI-based intrusion detection system '
                  'for encrypted traffic that integrates machine learning classification '
                  'with an interactive dashboard and multiple input modes.'],
         'right_title': 'Objectives',
         'right': ['Preprocess the CICIDS2017 dataset for training',
                   'Train and tune an XGBoost classifier on flow-based features',
                   'Build a Streamlit dashboard with Matrix-themed UI',
                   'Implement CSV upload, PCAP parsing, and live capture modes',
                   'Integrate an AI assistant for attack explanation',
                   'Evaluate model performance using standard metrics']},
        {'kicker': 'Slide 5', 'title': 'Significance of the Study',
         'items': [
             'Technological Impact: Demonstrates the practical application of gradient '
             'boosting for encrypted traffic analysis in a real-world tool.',
             'Economic Impact: Reduces the cost of network security monitoring by providing '
             'a free, open-source alternative to commercial IDS solutions.',
             'Educational Impact: Serves as a learning platform for understanding ML-based '
             'network security concepts and traffic analysis.',
             'Social Impact: Makes advanced network security accessible to small organisations '
             'and individuals who cannot afford expensive security infrastructure.',
         ]},
        {'kicker': 'Slide 6', 'title': 'Literature Review',
         'items': [
             'Buczak & Guven (2016) surveyed ML methods for cybersecurity, identifying '
             'random forests and neural networks as effective for intrusion detection.',
             'Chen & Guestrin (2016) introduced XGBoost, a scalable tree boosting system '
             'that became a benchmark for tabular data classification.',
             'Sharafaldin et al. (2018) created CICIDS2017, a realistic IDS dataset with '
             '71 flow features and 14 attack types, enabling reproducible research.',
             'Zhang et al. (2019) applied XGBoost to network intrusion detection, achieving '
             'high accuracy on multiple benchmark datasets.',
             'Gap: Few existing solutions combine high-accuracy ML with a modern interactive '
             'dashboard and multi-mode input (CSV, PCAP, live capture).',
         ]},
        {'kicker': 'Slide 7', 'title': 'Methodology',
         'kind': 'cards',
         'cards': [
             ('Approach', 'Agile SDLC with iterative development and continuous testing.'),
             ('Architecture', 'Streamlit frontend + XGBoost backend + tshark for capture.'),
             ('Dataset', 'CICIDS2017: 3.37M flows, 71 features, 14 attack types.'),
             ('Model', 'XGBoost classifier: 99.9% accuracy, 0.9999 ROC-AUC.'),
             ('Input Modes', 'CSV upload, PCAP parsing via tshark, live NIC capture.'),
             ('AI Assistant', 'Google Gemini for plain-English attack explanations.')]},
        {'kicker': 'Slide 8', 'title': 'Implementation',
         'kind': 'screenshot',
         'items': [
              ('Figure 1: Full Dashboard',
               'Matrix-themed IDS dashboard with hero section, feature cards, and attack tiles.',
               r'c:\Users\ALEXIS\Desktop\SENPAI\shots\ids\01-full-dashboard.png'),
              ('Figure 2: PCAP & Live Tabs',
               'PCAP upload interface and live network capture configuration.',
               r'c:\Users\ALEXIS\Desktop\SENPAI\shots\ids\02-pcap-tab.png'),
         ]},
        {'kicker': 'Slide 9', 'title': 'Results & Discussion',
         'kind': 'bullets',
         'items': [
              'Accuracy: 99.9% \u2014 999 out of 1000 flows correctly classified.',
              'ROC-AUC: 0.9999 \u2014 near-perfect discrimination between benign and malicious.',
              'Precision: 0.999 | Recall: 0.999 | F1-Score: 0.999.',
              'Confusion matrix shows near-diagonal dominance (minimal FP/FN).',
              'Feature importance: flow duration, backward packet length, inter-arrival times.',
              'All three modes (CSV, PCAP, Live) passed unit, integration, and UAT testing.',
              'AI assistant provides accurate plain-language explanations of predictions.',
         ]},
        {'kicker': 'Slide 10', 'title': 'Contributions to Knowledge',
         'items': [
             'A fully functional AI-based IDS achieving 99.9% accuracy and 0.9999 ROC-AUC '
             'on the CICIDS2017 benchmark dataset.',
             'Integration of three detection modes (CSV, PCAP, live) within a single '
             'unified dashboard accessible to non-experts.',
             'A Matrix-themed interactive interface that combines real-time detection '
             'with visual analytics and an AI-powered explanation system.',
             'Demonstration of XGBoost effectiveness for encrypted traffic analysis, '
             'with flow-based features requiring no decryption.',
             'A practical, open-source tool that bridges the gap between ML research '
             'and usable network security software.',
         ]},
        {'kicker': 'Slide 11', 'title': 'Conclusion / References',
         'kind': 'two',
         'left_title': 'Conclusion',
         'left': ['All six objectives were achieved: the CICIDS2017 dataset was '
                  'preprocessed, an XGBoost model was trained and tuned, a Streamlit '
                  'dashboard was built with the Matrix theme, CSV/PCAP/live modes were '
                  'implemented, a Gemini AI assistant was integrated, and the model '
                  'was evaluated achieving 99.9% accuracy and 0.9999 ROC-AUC.',
                  'The system provides a practical, accurate, and accessible solution '
                  'for encrypted traffic intrusion detection.'],
         'right_title': 'References',
         'right': ['Buczak, A. L. & Guven, E. (2016). IEEE COMST, 18(2), 1153\u20131176.',
                   'Chen, T. & Guestrin, C. (2016). KDD 2016, 785\u2013794.',
                   'Pressman, R. S. & Maxim, B. R. (2020). Software engineering (9th ed.). '
                   'McGraw-Hill.',
                   'Sharafaldin, I. et al. (2018). ICISSp, 108\u2013116.',
                   'Sommerville, I. (2016). Software engineering (10th ed.). Pearson.',
                   'Stallings, W. (2017). Network security essentials (6th ed.). Pearson.',
                   'Zhang, Y. et al. (2019). IEEE Access, 7, 164380\u2013164391.']},
    ],
}

build_deck(mayor, r'c:\Users\ALEXIS\Desktop\SENPAI\mayor4code_Defence_Slides.pptx')
build_deck(agri, r'c:\Users\ALEXIS\Desktop\SENPAI\AgriFlow_AI_Defence_Slides.pptx')
build_deck(agri_plain, r'c:\Users\ALEXIS\Desktop\SENPAI\AgriFlow_AI_Defence_Slides_Plain.pptx')
build_deck(ids, r'c:\Users\ALEXIS\Desktop\SENPAI\ids_Defence_Slides.pptx')
